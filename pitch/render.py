"""Render slides.html → PNGs → .pptx ready for Google Slides upload."""

import os
import glob
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches, Emu

SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
HTML_FILE = os.path.join(SCRIPT_DIR, "slides.html")
IMG_DIR = os.path.join(SCRIPT_DIR, "slide_images")
OUTPUT = os.path.join(SCRIPT_DIR, "z_experience_pitch.pptx")
SLIDE_COUNT = 21

# 16:9 at 2x
SLIDE_W = 1920
SLIDE_H = 1080


def render_pngs():
    os.makedirs(IMG_DIR, exist_ok=True)

    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(viewport={"width": SLIDE_W, "height": SLIDE_H})
        page.goto(f"file://{HTML_FILE}")
        # wait for fonts
        page.wait_for_timeout(2000)

        for i in range(1, SLIDE_COUNT + 1):
            el = page.locator(f"#slide-{i}")
            path = os.path.join(IMG_DIR, f"slide_{i:02d}.png")
            el.screenshot(path=path)
            print(f"  [{i:2d}/{SLIDE_COUNT}] → {path}")

        browser.close()


def build_pptx():
    prs = Presentation()
    # 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    images = sorted(glob.glob(os.path.join(IMG_DIR, "slide_*.png")))
    for img_path in images:
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        slide.shapes.add_picture(
            img_path,
            left=0,
            top=0,
            width=prs.slide_width,
            height=prs.slide_height,
        )
        print(f"  Added {os.path.basename(img_path)}")

    prs.save(OUTPUT)
    print(f"\nSaved → {OUTPUT}")


if __name__ == "__main__":
    print("Rendering PNGs from HTML...")
    render_pngs()
    print("\nBuilding .pptx...")
    build_pptx()
