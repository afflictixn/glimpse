"""Build the Z Experience pitch deck as .pptx (uploadable to Google Slides)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Palette ──────────────────────────────────────────────────
BG        = RGBColor(0x0A, 0x0A, 0x0A)
BG_DARK   = RGBColor(0x05, 0x05, 0x05)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_LIGHT = RGBColor(0xBB, 0xBB, 0xBB)
GRAY      = RGBColor(0x88, 0x88, 0x88)
GRAY_DIM  = RGBColor(0x55, 0x55, 0x55)
ACCENT    = RGBColor(0x6C, 0x63, 0xFF)  # purple
ACCENT2   = RGBColor(0x00, 0xD4, 0xAA)  # teal
ACCENT3   = RGBColor(0xFF, 0x6B, 0x6B)  # coral
FONT      = "Inter"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height


def set_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=32,
                 color=WHITE, bold=False, alignment=PP_ALIGN.CENTER,
                 font_name=FONT, line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(0)
    p.space_before = Pt(0)
    if line_spacing != 1.0:
        p.line_spacing = Pt(font_size * line_spacing)
    return txBox, tf


def add_para(tf, text, font_size=32, color=WHITE, bold=False,
             alignment=PP_ALIGN.CENTER, space_before=0, space_after=0,
             font_name=FONT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, corner=Inches(0.2)):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def center_x(box_w):
    return int((W - box_w) / 2)


# ═══════════════════════════════════════════════════════════════
# SLIDE 1 — "AI."
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(s, BG)
_, tf = add_text_box(s, center_x(Inches(6)), Inches(2), Inches(6), Inches(3.5),
                     "AI.", font_size=120, color=ACCENT, bold=True)

# ═══════════════════════════════════════════════════════════════
# SLIDE 2 — "What comes to mind?"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
_, tf = add_text_box(s, center_x(Inches(10)), Inches(2.2), Inches(10), Inches(3),
                     "What comes to mind", font_size=56, color=GRAY, bold=False)
add_para(tf, "when you see this word?", font_size=56, color=GRAY)

# ═══════════════════════════════════════════════════════════════
# SLIDE 3 — "Probably this." (three chat mockups)
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
add_text_box(s, center_x(Inches(6)), Inches(0.5), Inches(6), Inches(1),
             "Probably this.", font_size=36, color=GRAY_DIM, bold=False)

chats = [
    ("ChatGPT", "Write me an essay\nabout climate change",
     "Climate change is one of the\nmost pressing issues facing..."),
    ("Claude", "Help me debug this\nPython error",
     "I'd be happy to help!\nCould you share the traceback..."),
    ("Gemini", "Summarize this article\nfor me",
     "Here are the key takeaways\nfrom the article..."),
]

card_w = Inches(3.5)
card_h = Inches(4.2)
gap = Inches(0.5)
total_w = card_w * 3 + gap * 2
start_x = int((W - total_w) / 2)

for i, (name, user_msg, ai_msg) in enumerate(chats):
    x = start_x + i * (card_w + gap)
    y = Inches(1.8)

    # Card background
    add_rounded_rect(s, x, y, card_w, card_h, RGBColor(0x1A, 0x1A, 0x1A))

    # Title
    add_text_box(s, x + Inches(0.3), y + Inches(0.25), Inches(2.5), Inches(0.4),
                 name.upper(), font_size=11, color=GRAY_DIM, bold=True,
                 alignment=PP_ALIGN.LEFT)

    # User bubble
    ub = add_rounded_rect(s, x + Inches(0.8), y + Inches(0.9), Inches(2.4), Inches(1.0), ACCENT)
    add_text_box(s, x + Inches(0.9), y + Inches(0.95), Inches(2.2), Inches(0.9),
                 user_msg, font_size=13, color=WHITE, alignment=PP_ALIGN.LEFT)

    # AI bubble
    ab = add_rounded_rect(s, x + Inches(0.3), y + Inches(2.1), Inches(2.4), Inches(1.0),
                          RGBColor(0x25, 0x25, 0x25))
    add_text_box(s, x + Inches(0.4), y + Inches(2.15), Inches(2.2), Inches(0.9),
                 ai_msg, font_size=13, color=GRAY_LIGHT, alignment=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════
# SLIDE 4 — "Who here imagines AI like this?"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
add_text_box(s, center_x(Inches(10)), Inches(2.5), Inches(10), Inches(2),
             "Who here imagines AI like this?", font_size=52, color=WHITE, bold=False)
add_text_box(s, center_x(Inches(4)), Inches(4.5), Inches(4), Inches(1),
             "\U0001F64B \U0001F64B\u200D\u2642\uFE0F \U0001F64B\u200D\u2640\uFE0F",
             font_size=48, color=WHITE)

# ═══════════════════════════════════════════════════════════════
# SLIDE 5 — The problem with chat
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
add_text_box(s, center_x(Inches(10)), Inches(0.5), Inches(10), Inches(1),
             "The problem with chat.", font_size=44, color=WHITE, bold=True)

problems = [
    ("\U0001F4DD", "You type the entire context. Every. Single. Time.",
     "Copy-paste your code, explain the situation, set the scene..."),
    ("\U0001F9E0", '"When did we talk about that?"',
     "Context is lost between sessions. You start from scratch."),
    ("\u23F3", "You have to ask first.",
     "AI just sits there. Waiting. It never comes to you."),
]

for i, (icon, title, desc) in enumerate(problems):
    y = Inches(1.8) + i * Inches(1.7)
    add_text_box(s, Inches(2.5), y, Inches(0.8), Inches(0.8),
                 icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_text_box(s, Inches(3.5), y, Inches(7), Inches(0.5),
                 title, font_size=24, color=WHITE, bold=True,
                 alignment=PP_ALIGN.LEFT)
    add_text_box(s, Inches(3.5), y + Inches(0.5), Inches(7), Inches(0.5),
                 desc, font_size=18, color=GRAY,
                 alignment=PP_ALIGN.LEFT)

# "Raise your hand" prompt
add_text_box(s, center_x(Inches(8)), Inches(6.2), Inches(8), Inches(0.6),
             "Sound familiar?", font_size=20, color=GRAY_DIM)

# ═══════════════════════════════════════════════════════════════
# SLIDE 6 — "Doesn't feel quite REVOLUTIONARY, right?"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
_, tf = add_text_box(s, center_x(Inches(10)), Inches(2), Inches(10), Inches(1.2),
                     "Doesn't feel quite", font_size=52, color=GRAY_DIM)
# "revolutionary" with strikethrough
p = add_para(tf, "revolutionary", font_size=52, color=ACCENT3, bold=True, space_before=8)
p.font.strikethrough = True
add_para(tf, "right?", font_size=52, color=GRAY_DIM, space_before=8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 7 — "It doesn't have to be like this."
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
_, tf = add_text_box(s, center_x(Inches(10)), Inches(2.5), Inches(10), Inches(2.5),
                     "It doesn't have to", font_size=56, color=GRAY)
add_para(tf, "be like this.", font_size=56, color=GRAY, space_before=8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 8 — "Introducing"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text_box(s, center_x(Inches(6)), Inches(3), Inches(6), Inches(1.5),
             "INTRODUCING", font_size=24, color=GRAY_DIM, bold=False)

# ═══════════════════════════════════════════════════════════════
# SLIDE 9 — "Z" hero
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
add_text_box(s, center_x(Inches(6)), Inches(1.5), Inches(6), Inches(3),
             "Z", font_size=180, color=ACCENT, bold=True)
add_text_box(s, center_x(Inches(6)), Inches(4.8), Inches(6), Inches(1),
             "E X P E R I E N C E", font_size=22, color=GRAY_DIM)

# ═══════════════════════════════════════════════════════════════
# SLIDE 10 — "A new experience of using your computer."
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)
_, tf = add_text_box(s, center_x(Inches(10)), Inches(2.5), Inches(10), Inches(2.5),
                     "A new experience", font_size=52, color=GRAY)
add_para(tf, "of using your computer.", font_size=52, color=WHITE, bold=True, space_before=8)

# ═══════════════════════════════════════════════════════════════
# SLIDE 11 — Three pillars: Proactive / Caring / Private
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG)

pillars = [
    ("\u26A1", "PROACTIVE", ACCENT,
     "Z doesn't wait for you to ask.\nIt sees what you're doing and\noffers help before you need it."),
    ("\U0001F49C", "CARING", ACCENT2,
     "Z remembers your context,\nyour conversations, your habits.\nIt genuinely looks out for you."),
    ("\U0001F512", "PRIVATE", ACCENT3,
     "Everything runs locally on\nyour machine. Your data never\nleaves your computer."),
]

col_w = Inches(3.5)
col_gap = Inches(0.4)
total = col_w * 3 + col_gap * 2
sx = int((W - total) / 2)

for i, (icon, title, clr, desc) in enumerate(pillars):
    x = sx + i * (col_w + col_gap)
    add_text_box(s, x, Inches(1.5), col_w, Inches(1),
                 icon, font_size=52, alignment=PP_ALIGN.CENTER)
    add_text_box(s, x, Inches(2.7), col_w, Inches(0.6),
                 title, font_size=22, color=clr, bold=True)
    add_text_box(s, x, Inches(3.5), col_w, Inches(2),
                 desc, font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER,
                 line_spacing=1.5)

# ═══════════════════════════════════════════════════════════════
# SLIDE 12 — "Z, how about you just introduce yourself?"
# ═══════════════════════════════════════════════════════════════
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, BG_DARK)
_, tf = add_text_box(s, center_x(Inches(10)), Inches(2), Inches(10), Inches(2.5),
                     "Z,", font_size=52, color=ACCENT, bold=True)
add_para(tf, "how about you just", font_size=52, color=GRAY, space_before=8)
add_para(tf, "introduce yourself?", font_size=52, color=WHITE, bold=True, space_before=8)

# Live demo badge
badge = add_rounded_rect(s, center_x(Inches(2.5)), Inches(5.5), Inches(2.5), Inches(0.7),
                         BG_DARK)
badge.line.color.rgb = ACCENT
badge.line.width = Pt(2)
add_text_box(s, center_x(Inches(2.5)), Inches(5.55), Inches(2.5), Inches(0.6),
             "LIVE DEMO", font_size=16, color=ACCENT, bold=True)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
out = os.path.join(os.path.dirname(__file__), "z_experience_pitch.pptx")
prs.save(out)
print(f"Saved → {out}")
